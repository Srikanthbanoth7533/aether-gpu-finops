import os
import sys
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Define slide layout constants (index 6 is usually a blank layout)
BLANK_LAYOUT_INDEX = 6

# Colors (Zinc Dark Theme Palette)
BG_COLOR = RGBColor(12, 12, 15)      # #0c0c0f
TEXT_MAIN = RGBColor(250, 250, 250)  # #fafafa
TEXT_MUTED = RGBColor(161, 161, 170) # #a1a1aa
PRIMARY = RGBColor(59, 130, 246)     # #3b82f6
SUCCESS = RGBColor(16, 185, 129)     # #10b981
WARNING = RGBColor(245, 158, 11)     # #f59e0b
DANGER = RGBColor(239, 68, 68)       # #ef4444

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_slide_title(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    return txBox

def add_slide_body(slide, bullet_points):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    for i, bp in enumerate(bullet_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(14)
        
        # Split on first colon to bold the title
        if ":" in bp:
            title, desc = bp.split(":", 1)
            run_title = p.add_run()
            run_title.text = title + ":"
            run_title.font.name = 'Arial'
            run_title.font.size = Pt(18)
            run_title.font.bold = True
            run_title.font.color.rgb = TEXT_MAIN
            
            run_desc = p.add_run()
            run_desc.text = desc
            run_desc.font.name = 'Arial'
            run_desc.font.size = Pt(18)
            run_desc.font.color.rgb = TEXT_MUTED
        else:
            run = p.add_run()
            run.text = bp
            run.font.name = 'Arial'
            run.font.size = Pt(18)
            run.font.color.rgb = TEXT_MAIN

def add_slide_body_with_image(slide, bullet_points, image_path):
    # Left column: Text (narrower width to accommodate the image)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(3.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    for i, bp in enumerate(bullet_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(10)
        
        # Split on first colon to bold the title
        if ":" in bp:
            title, desc = bp.split(":", 1)
            run_title = p.add_run()
            run_title.text = title + ":"
            run_title.font.name = 'Arial'
            run_title.font.size = Pt(13)
            run_title.font.bold = True
            run_title.font.color.rgb = TEXT_MAIN
            
            run_desc = p.add_run()
            run_desc.text = desc
            run_desc.font.name = 'Arial'
            run_desc.font.size = Pt(13)
            run_desc.font.color.rgb = TEXT_MUTED
        else:
            run = p.add_run()
            run.text = bp
            run.font.name = 'Arial'
            run.font.size = Pt(13)
            run.font.color.rgb = TEXT_MAIN
            
    # Right column: Image
    if os.path.exists(image_path):
        # Position image on the right side
        slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.4), width=Inches(4.3))
    else:
        print(f"WARNING: Image not found at {os.path.abspath(image_path)}")

def create_presentation():
    prs = pptx.Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625) # 16:9 aspect ratio
    blank_layout = prs.slide_layouts[BLANK_LAYOUT_INDEX]
    
    # -------------------------------------------------------------------------
    # Slide 1: Title Slide
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "AETHERFIN GPU OPA"
    p1.text = "AETHERFIN GPU OPS"
    p1.font.name = 'Arial'
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p1.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "GenAI GPU FinOps & Carbon Sustainability Platform"
    p2.font.name = 'Arial'
    p2.font.size = Pt(22)
    p2.font.color.rgb = SUCCESS
    p2.space_after = Pt(30)
    
    p3 = tf.add_paragraph()
    p3.text = "Developed by: Banoth Srikanth\nPortfolio Deliverable for Enterprise Product Analyst & Data Engineer Reviews"
    p3.font.name = 'Arial'
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_MUTED
    
    # -------------------------------------------------------------------------
    # Slide 2: Executive Summary
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Executive Summary")
    add_slide_body(slide, [
        "Core Value Proposition: Bridges the visibility gap between raw machine load telemetry and business-level finance/environmental indexes.",
        "Total Compute Spend: Realized a baseline operational spend of $7,844.98 across 5 active cluster nodes.",
        "GICP Cost Waste Penalty: Identified $5,073.76 (64.7%) in wasted spend on idle instances running under 15% load.",
        "Compute Efficiency: Serving average of 2.23 tokens per Wh, benchmarked across model deployments.",
        "Environmental Impact: Tracked 467.32 kg in Scope 3 carbon footprint emissions."
    ])

    # -------------------------------------------------------------------------
    # Slide 3: Problem Statement
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Problem Statement")
    add_slide_body(slide, [
        "Skyrocketing Infrastructure Budgets: Generative AI model training and serving represent the largest modern enterprise IT expense.",
        "The Visibility Vacuum: Standard server tools show CPU/Memory load but omit cost, SLA penalties, and carbon metrics.",
        "Severe Resource Waste: Multi-million dollar GPU clusters remain reserved 24/7 with zero load outside working hours.",
        "Regulatory Reporting Deficits: Compliance audits require tracking Scope 3 green IT emissions, but standard cloud metrics fail to correlate electricity grid coefficients."
    ])

    # -------------------------------------------------------------------------
    # Slide 4: Business Context
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Business Context")
    add_slide_body(slide, [
        "Financial Leakage: Over 60% of GPU compute budgets are wasted on idle allocations.",
        "Regulatory Compliance: Imminent Scope 3 environmental carbon footprint audits are legally required in global jurisdictions.",
        "The Green IT Opportunity: Optimizing compute sizing and regional scheduling yields double benefits: lowering spend and carbon footprint.",
        "Stakeholder Alignment: Designed to support Chief Financial Officers (CFOs), VP of Engineering, AI Product Managers, and Sustainability Officers."
    ])

    # -------------------------------------------------------------------------
    # Slide 5: Product Vision
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Product Vision")
    add_slide_body(slide, [
        "Continuous Intelligence: Convert low-level machine logs into high-level business analytics indexes.",
        "Actionable Sizing Engine: Feed recommendation pipelines directly to container orchestration systems to shut down waste.",
        "Risk-Balanced Simulation: Allow product managers to simulate cost and SLA performance changes before allocating clusters.",
        "Immutable Audit Trails: Guarantee security compliance through JWT auth and immutable transaction auditing logs."
    ])

    # -------------------------------------------------------------------------
    # Slide 6: System Architecture
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "System Architecture")
    add_slide_body(slide, [
        "Multi-Tier Design: Decoupled client, REST API server, and relational database tiers.",
        "Frontend Client: React TypeScript application executing custom ECharts rendering inside a premium Glassmorphism Zinc dark UI.",
        "FastAPI Backend: High-performance Python ASGI API routing, using NumPy for sliding-window mathematical Z-scores.",
        "Database Layer: SQLAlchemy ORM layer abstracting SQLite (development) and MySQL (production) servers.",
        "Orchestration Stacks: Dockerized compose profiles managing application services, database engines, and pgAdmin monitors."
    ])

    # -------------------------------------------------------------------------
    # Slide 7: Database Design
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Database Design & Schema")
    add_slide_body(slide, [
        "gpu_cluster_nodes: Node-level registry tracking provider (AWS/GCP/On-Prem), region, instance type, hourly cost, and grid carbon intensity.",
        "model_deployments: Active serving configurations mapping allocated GPUs, latency SLAs, and target TPS thresholds.",
        "gpu_utilization_logs: Telemetry logging hourly average GPU loads and actual power draw in Watts.",
        "carbon_emissions_logs: Correlating power draws to regional grid intensity to log emissions hourly.",
        "inference_request_logs: Transaction logging query parameters, token numbers, latency, and SLA violation status."
    ])

    # -------------------------------------------------------------------------
    # Slide 8: Technology Stack
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Technology Stack")
    add_slide_body(slide, [
        "Frontend: React 18, TypeScript, Vite, ECharts-for-React, Lucide icons, and custom CSS.",
        "Backend: Python 3.11+, FastAPI, Uvicorn, SQLAlchemy, PyJWT, Cryptography, and NumPy.",
        "Database: MySQL (Local/Production), SQLite (Testing), PyMySQL, and PostgreSQL support.",
        "Deployment & DevOps: Docker, Docker Compose, Git repository source control, and Vercel hosting integration."
    ])

    # -------------------------------------------------------------------------
    # Slide 9: KPI Framework
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Innovative KPI Framework")
    add_slide_body(slide, [
        "GPU Idle Cost Penalty (GICP): Translates technical load metrics into dollar waste assessments.",
        "Model Compute Efficiency Score (MCES): Tracks token output density relative to physical energy consumption.",
        "SLA Violation Exposure Index (SLA-VEI): Measures user latency breaches, weighting extreme delays non-linearly.",
        "Carbon Offset ROI (COROI): Correlates financial savings of regional rerouting actions against offset expenses."
    ])

    # -------------------------------------------------------------------------
    # Slide 10: GICP Explanation
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "GPU Idle Cost Penalty (GICP)")
    add_slide_body(slide, [
        "Formula: GICP = sum(Idle Hours * Hourly Cost) / Total Spend * 100",
        "Threshold Definition: Tracks hours where a node is active but average GPU load stays under 15%.",
        "Business Cost: In the audited trial, $5,073.76 out of $7,844.98 was wasted on idle capacity.",
        "Outcome Action: Identifies exact instances suitable for decommissioning or dynamic scale-down."
    ])

    # -------------------------------------------------------------------------
    # Slide 11: MCES Explanation
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Model Compute Efficiency Score (MCES)")
    add_slide_body(slide, [
        "Formula: MCES = (Prompt Tokens + Completion Tokens) / sum(Power Draw in Watts * 1 Hour)",
        "Standardization Value: Allows C-suite comparison of model architectures on different hardware types.",
        "Hardware Insights: Llama-3-8B on GCP L4 nodes achieved 8.42 tokens/Wh, compared to 1.12 tokens/Wh for Llama-3-70B on H100s.",
        "Outcome Action: Identifies candidate models for right-sizing or quantization."
    ])

    # -------------------------------------------------------------------------
    # Slide 12: SLA-VEI Explanation
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "SLA Violation Exposure Index (SLA-VEI)")
    add_slide_body(slide, [
        "Formula: SLA-VEI = sum(max(0, (Actual Latency - Target Latency) / Target Latency)) / Total Requests * 100",
        "The Severity Weighting: Counts of violations do not capture user frustration; SLA-VEI weights delays exponentially.",
        "Baseline Score: Current baseline score is 2.89, driven by peak diurnal traffic queuing delays.",
        "Outcome Action: Signals when cost-cutting measures threaten customer retention."
    ])

    # -------------------------------------------------------------------------
    # Slide 13: COROI Explanation
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Carbon Offset ROI (COROI)")
    add_slide_body(slide, [
        "Index Concept: Evaluates the cost efficiency of shifting workloads to low-intensity regional grids.",
        "Grid Differences: Compares On-Premises coal grid (650g CO2/kWh) to GCP Europe-West1 nuclear grid (45g CO2/kWh).",
        "The Trade-Off: Carbon routing may increase network latencies but completely eliminates carbon penalty offsets.",
        "Outcome Action: Scheduler routes batch jobs to clean regions during off-peak times."
    ])

    # -------------------------------------------------------------------------
    # Slide 14: Dashboard Screenshots [WITH IMAGE]
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Executive & Dashboard Previews")
    add_slide_body_with_image(slide, [
        "Dashboard (sizing.png): Features the high-level KPI overview panels showing spend and carbon footprint details.",
        "Compute Activity: Renders the active 48-hour timeline detailing node utilization and load fluctuations.",
        "Node Inventory: Renders active database logs of multi-cloud allocations.",
        "One-Click Exports: Stream summaries and audit trails directly to CSV or formatted text reports."
    ], "../screenshots/sizing.png")

    # -------------------------------------------------------------------------
    # Slide 15: Trend Analysis [WITH IMAGE]
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Trend & Cohort Analytics")
    add_slide_body_with_image(slide, [
        "Interactive Sidebar (drilldown.png): Opens node-level financial details dynamically on click.",
        "Visual Sizing Factors: CFOs can drill down to see exact provider and regional billing components.",
        "Multi-Range Filters: Supports 7-day, 30-day, and custom timeline aggregation windows.",
        "Cohort Analysis: Compares model execution profiles to pinpoint underperforming instances."
    ], "../screenshots/drilldown.png")

    # -------------------------------------------------------------------------
    # Slide 16: Anomaly Detection [WITH IMAGE]
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Z-Score Anomaly Detection")
    add_slide_body_with_image(slide, [
        "Outliers Graph (anomalies.png): Renders Z-score statistical outlier plots tracking latency spikes.",
        "Waste Alerts: Automatically flags nodes operating under 10% average load for 12+ consecutive hours.",
        "Root Cause Diagnostics: Automatically pairs anomalies with technical diagnoses and routing solutions.",
        "Dynamic Baselines: Avoids static thresholds by computing sliding-window mean and standard deviations."
    ], "../screenshots/anomalies.png")

    # -------------------------------------------------------------------------
    # Slide 17: AI Advisor [WITH IMAGE]
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "AI FinOps Advisor")
    add_slide_body_with_image(slide, [
        "Conversational Interface (casestudy.png): Provides context-enriched advisory panels displaying strategic plans.",
        "Direct Answers: Responds to inquiries like 'How can I reduce costs by 20%?' with detailed, multi-step actions.",
        "Fallback Rules: Implements rule-based fallback logic to analyze local database statistics if the LLM API is offline.",
        "Action Chips: Interactive cards allowing users to click common queries instantly."
    ], "../screenshots/casestudy.png")

    # -------------------------------------------------------------------------
    # Slide 18: What-If Simulator
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "What-If Sizing Simulator")
    add_slide_body(slide, [
        "Predictive Sizing Controls: Adjust scale multiplier, daily operating hours, and traffic concurrency.",
        "Routing Strategies: Toggle default, Cost-Optimized (cheapest), and Carbon-Optimized (cleanest) modes.",
        "Trade-Off Modeling: Simulates how cost-cutting changes affect SLA violation risks.",
        "Real-Time Previews: Computes projected costs, carbon footprints, and net savings dynamically."
    ])

    # -------------------------------------------------------------------------
    # Slide 19: Load Testing Results
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Performance & Load Audits")
    add_slide_body(slide, [
        "Optimization Focus: Resolved DB query loops by pre-fetching deployment records into memory dictionaries.",
        "100 Concurrency: 100% Success | P50 latency 804.1ms | Peak RAM 93.9 MB",
        "500 Concurrency: 100% Success | P50 latency 886.2ms | Peak RAM 95.3 MB",
        "1,000 Concurrency: 100% Success | P50 latency 967.6ms | Peak RAM 96.9 MB"
    ])

    # -------------------------------------------------------------------------
    # Slide 20: Security Validation [WITH IMAGE]
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Security & Compliance Audits")
    add_slide_body_with_image(slide, [
        "Audit Log Viewer (admin_logs.png): Renders secure table logging details of all system modifications.",
        "Access Protection: Bearer JWT validation shields all secure API routes and telemetry reports.",
        "Type Safety Validation: Pydantic schemas filter and reject malformed JSON payloads (HTTP 422).",
        "SQL Injection Prevention: Parameterized queries in SQLAlchemy ORM eliminate injection vectors."
    ], "../screenshots/admin_logs.png")

    # -------------------------------------------------------------------------
    # Slide 21: Business Impact
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Business Impact Summary")
    add_slide_body(slide, [
        "Double-Bottom Line Value: Shows that optimizing compute resources reduces spend and carbon emissions simultaneously.",
        "Cost Savings: Cut weekly cluster run-rates by 50% through smart sizing.",
        "Carbon Reduction: Achieved a 73% drop in Scope 3 greenhouse gas footprints.",
        "Enterprise Readiness: Audit trails and RBAC controls satisfy compliance audits."
    ])

    # -------------------------------------------------------------------------
    # Slide 22: Cost Savings
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Cost Optimization Actions")
    add_slide_body(slide, [
        "Terminating Idle Node: Decommissioning `aws-us-east-1-a100-idle` saves $2,284.80 / week.",
        "Off-Peak Auto-Scaling: Downsizing H100 counts by 50% between 11 PM and 6 AM saves $478.24 / week.",
        "Model Right-Sizing: Migrating Mixtral from GCP A100 to GCP L4 instances saves $806.40 / week.",
        "Net Sizing Business Outcome: Delivers $3,921.58 / week ($16,980 monthly savings, 50.0% cost reduction)."
    ])

    # -------------------------------------------------------------------------
    # Slide 23: Carbon Reduction
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Environmental Sustainability Outcomes")
    add_slide_body(slide, [
        "Grid Schedulers: Batch jobs are automatically routed to cleaner regional grids during off-peak hours.",
        "Stable-Diffusion XL Migration: Shifting image batch jobs to GCP Europe West-1 wind grid saves 343 kg CO2 / week.",
        "Net Environmental Sizing: Lowers weekly emissions from 467.3 kg CO2 to 124.2 kg CO2 (73.4% reduction).",
        "Offsets Optimization: Avoids $356 / month in carbon offsets, accelerating ESG timeline targets."
    ])

    # -------------------------------------------------------------------------
    # Slide 24: Future Roadmap
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Future Roadmap")
    add_slide_body(slide, [
        "Autonomous Cloud Scaling: Connect the sizing API to AWS/GCP dynamic node groups to scale pools directly.",
        "Predictive Load Forecasting: Implement regression models to forecast peak traffic and pre-allocate spot instances.",
        "Direct Offset Marketplaces: Integrate carbon credit purchasing APIs directly into the management panel.",
        "Multi-Tenant SaaS: Add organization isolation schemas to serve multiple enterprise accounts."
    ])

    # -------------------------------------------------------------------------
    # Slide 25: Conclusion [WITH IMAGE]
    # -------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_slide_title(slide, "Conclusion")
    add_slide_body_with_image(slide, [
        "Recruiter-Ready Platform: Successfully bridges GenAI compute operations, financial auditing, and carbon ESG metrics.",
        "Designed UI (login.png): Offers secure credentials authentication layout using Zinc dark themes.",
        "Production-Grade Performance: Validated at 1,000 requests concurrency under 100% success rates.",
        "Advanced Competency: Demonstrates full capability in Product Analytics, Data Engineering, and Security."
    ], "../screenshots/login.png")

    output_path = "../aetherfin_presentation.pptx"
    prs.save(output_path)
    print(f"PowerPoint presentation generated at: {os.path.abspath(output_path)}")
    return os.path.abspath(output_path)

def convert_pptx_to_pdf(pptx_path):
    pdf_path = pptx_path.replace(".pptx", ".pdf")
    print(f"Converting {pptx_path} to {pdf_path}...")
    
    # Initialize COM automation
    import win32com.client
    
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1
    
    try:
        deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
        # Format Type 32 is PDF format in PowerPoint COM
        deck.SaveAs(os.path.abspath(pdf_path), 32)
        deck.Close()
        print("PDF conversion completed successfully!")
        return os.path.abspath(pdf_path)
    except Exception as e:
        print(f"ERROR during PDF conversion: {e}")
        sys.exit(1)
    finally:
        powerpoint.Quit()

if __name__ == "__main__":
    pptx_file = create_presentation()
    convert_pptx_to_pdf(pptx_file)
